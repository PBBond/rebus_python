import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import random
import mysql.connector
from mysql.connector import Error
import math
from os.path import exists as file_exists, basename
from flask import Flask, render_template, request, send_file
import requests
import json
import langid



#word = input("Enter some word for the puzzle: ")
#word = "అలంకరణ"
englishTeluguWordList = []


def getChars(input_str, language='English'):
    ws_api = "https://indic-wp.thisisjava.com/api/getLogicalChars.php"
    params = {"string": input_str, "language": language}

    # On some servers, Apache may reject the crawlers.
    # To mimic an actual user, send this dummy header along
    # https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/User-Agent
    headers = {'User-Agent': ''}

    # get the response object
    response_obj = requests.get(url=ws_api, params=params, headers=headers)

    # get the json response
    json_response = response_obj.text

    # Get rid of UTF-8 Byte Order Mark (BOM)
    if json_response.startswith(u'\ufeff'):
        json_response = json_response.encode('utf8')[6:].decode('utf8')

    # Load the json response to convert it a dictionary
    json_dict = json.loads(json_response)

    # Get the logical characters (spaces are also counted)
    logical_characters = json_dict['data']

    # return the list
    return logical_characters

#pw = input("Enter the database password: ")
pw = 'vasya316'
db = 'rebus'
#connection = createServerConnection("localhost", "root", pw)


def createDBConnection(hostName, userName, userPassword, dbName):
    connection = None
    try:
        connection = mysql.connector.connect(host = hostName, user = userName, password = userPassword, db = dbName)
        print("MySQL database connection successful")
    except Error as e:
        print(f'Error: {e}')
    return connection

def readQuery(connection, query):
    cursor = connection.cursor()
    result = None
    try:
        cursor.execute(query)
        result = cursor.fetchall()
        return result
    except Error as e:
        print(f"{e}")

q1 = "select * from words;"
connection = createDBConnection("localhost", "root", pw, db)
result = readQuery(connection, q1)



for row in result:
    row = list(row)
    if row[3] == '':
        row[3] = 'noImage.jpg'
    tempDict = {
       'telugu': row[1],
       'word': row[2],
       'image': f'http://rebus.telugupuzzles.com/Images/{row[3]}',
    }
    englishTeluguWordList.append(tempDict)

def removeExclusionWords(wordList, exclusionList):
    for item in wordList:
        for ex in exclusionList:
            if ex in item['word'] or ex in item['telugu']:
                wordList.remove(item)
    return wordList

def getWordListEnglish(givenWord, exclusion=None):
    tempList = englishTeluguWordList.copy()
    if exclusion:
        tempList = removeExclusionWords(tempList, exclusion)
    random.shuffle(tempList)
    wordDict =[]
    #givenWord = givenWord.replace(' ', '')
    for char in givenWord:
        for item in tempList:
            word = item['word']
            if char in word:
                l = [word, word.index(char), item['image']]
                wordDict.append(l)
                tempList.remove(item)
                break
    return wordDict

def getWordListTelugu(givenWord, exclusion=None):
    tempList = englishTeluguWordList.copy()
    if exclusion:
        tempList = removeExclusionWords(tempList, exclusion)
    random.shuffle(tempList)
    wordDict =[]
    #givenWord = givenWord.replace(' ', '')
    for char in givenWord:
        for item in tempList:
            word = item['telugu']
            if char in word:
                l = [word, word.index(char), item['image']]
                wordDict.append(l)
                tempList.remove(item)
                break
    return wordDict

def makeSlide(pr1, puzzleNum, language, logicalWord, showAns, dist = 1.5, type = 'default', exclusion = None):
    slide = pr1.slides.add_slide(pr1.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(2), Inches(0.2), Inches(5), Inches(1))
    tf = title.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Puzzle #{puzzleNum}"

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(64)
    font.bold = True

    numRows = math.ceil(len(logicalWord) / 4)
    numCols = 4

    width, height = Inches(dist), Inches(dist)
    if type == 'default':
        width, height = Inches(1.5), Inches(1.5)
    elif type == 'width':
        width = Inches(dist)
        if dist > 2.5:
            numCols -= 1
            numRows = math.ceil(len(logicalWord) / numCols)
    elif type == 'height':
        height = Inches(dist)
        if dist > 2.5:
            numRows += 1
            numCols -= 1
    # might need this tempWord = [char for char in word]
    #logicalWord = getChars(word)
    list_of_words = []
    if language == 'te':
        list_of_words = getWordListTelugu(logicalWord, exclusion)
    elif language == 'en':
        list_of_words = getWordListEnglish(logicalWord, exclusion)


    print(logicalWord)

    toReturn = list_of_words.copy()
    textColor = RGBColor(255, 120, 210)

    for j in range(numRows):
        topPic = Inches((j * 2) + 1.5)
        topWord = Inches((j * 2) + 3)
        for i in range(numCols):
            if not list_of_words:
                break
            #print(f'{list_of_words[0][2]}')
            pictureURL = list_of_words[0][2]
            r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
            if r.status_code == 200:
                # print(pictureURL)
                try:
                    with open(basename(pictureURL), "wb") as f:
                        r.raw.decode_content = True
                        shutil.copyfileobj(r.raw, f)
                except:
                    pass
            if type == 'default':
                try:
                    pic = slide.shapes.add_picture(basename(list_of_words[0][2]), Inches(1 + (i*2)), topPic, width=width, height=height)
                except:
                    pic = slide.shapes.add_picture(f'static/images/_not_found.png', Inches(1 + (i * 2)), topPic,
                                                   width=width, height=height)
                tb = slide.shapes.add_textbox(Inches(1 + (i*2)), topWord, Inches(1), Inches(0.5))

                tb.text = f'{list_of_words[0][1] + 1}/{len(list_of_words[0][0])}' # add the actual word {list_of_words[0][0]}
                list_of_words.pop(0)
            elif type == 'width':
                try:
                    pic = slide.shapes.add_picture(basename(list_of_words[0][2]), Inches(1 + (i * 2)), topPic,
                                                   width=width)
                except:
                    pic = slide.shapes.add_picture(f'static/images/_not_found.png', Inches(1 + (i * 2)), topPic,
                                                   width=width, height=height)
                tb = slide.shapes.add_textbox(Inches(1 + (i * 2)), topPic, Inches(1), Inches(0.5))
                tf = tb.text_frame

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f'{list_of_words[0][1] + 1}/{len(list_of_words[0][0])}'

                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)
                font.bold = True
                font.color.rgb = textColor
                list_of_words.pop(0)
            elif type == 'height':
                try:
                    pic = slide.shapes.add_picture(basename(list_of_words[0][2]), Inches(1 + (i * 2)), topPic,
                                                   width=width)
                except:
                    pic = slide.shapes.add_picture(f'static/images/_not_found.png', Inches(1 + (i * 2)), topPic, height=height)
                tb = slide.shapes.add_textbox(Inches(1 + (i * 2)), topPic, Inches(1), Inches(0.5))
                tf = tb.text_frame

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f'{list_of_words[0][1] + 1}/{len(list_of_words[0][0])}'

                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)
                font.bold = True
                font.color.rgb = textColor
                list_of_words.pop(0)

    return toReturn


def makeAnswerSlides(puzzles, pr1):
    # word, index, pic
    strings = []
    i = 1
    for puzzle in puzzles:
        string = f'#{i}'
        i += 1
        for item in puzzle:
            string += f' ({item[1] + 1}/{len(item[0])}) {item[0]} |'
        strings.append(string[:-2])

    #split up the powerpoint pages into max n items.
    n = 8
    chunks = [strings[i:i + n] for i in range(0, len(strings), n)]

    for chunk in chunks:
        Layout = pr1.slide_layouts[6]
        slide = pr1.slides.add_slide(Layout)
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        textframe = textbox.text_frame
        for string in chunk:
            para = textframe.add_paragraph()
            para.text = str(string)
        textframe.fit_text()


def getManyLists(searchWord, language, amount, exclusion=None):
    result = []
    if language =='te':
        for i in range(amount):
            tempList = getWordListTelugu(searchWord, exclusion)
            result.append(tempList)
    elif language =='en':
        for i in range(amount):
            tempList = getWordListEnglish(searchWord, exclusion)
            result.append(tempList)
    return result


app = Flask(__name__)

@app.route('/')
@app.route('/home')
def index():
    return render_template("index.html")

@app.route('/oneWordMany', methods=['POST', 'GET'])
def oneWordMany():
    if request.method == 'POST':
        puzzle_word = request.form['puzzle_word']
        logicalWord = getChars(puzzle_word)
        amount = request.form['owmp_amount']
        amount = int(amount)

        print(amount)
        #print(f"logical chars: {type(langid.classify(puzzle_word))}")

        allPuzzles = []
        if langid.classify(puzzle_word)[0] == 'te':
            allPuzzles = getManyLists(logicalWord, 'te', amount)
        else:
            allPuzzles = getManyLists(logicalWord, 'en', amount)
        return render_template('oneWordMany.html', load=True, puzzle_word=puzzle_word, all_puzzles=allPuzzles)
    else:
        return render_template('oneWordMany.html', load=False)




def makePowerPoint(language, logicalWord, showAnswers, amount, type, size):
    pr1 = Presentation()
    #hide the answers
    puzzles = []
    #if word is greater than 12 letters, it's split up amongst multiple pages
    n = 12
    chunks = [logicalWord[i:i + n] for i in range(0, len(logicalWord), n)]

    for i in range(amount):
        for chunk in chunks:
            puzzles.append(makeSlide(pr1, (i + 1), language, chunk, False, size, type))
    if showAnswers:
        makeAnswerSlides(puzzles, pr1)
    pr1.save('Rebus.pptx')

@app.route('/return-file')
def return_file():
    return send_file()

@app.route('/oneWordManyPPT', methods=['POST', 'GET'])
def oneWordManyPPT():
    if request.method == 'POST':
        size = float(request.form['size_value'])
        type = request.form['image_size']
        puzzle_word = request.form['puzzle_word']
        logicalWord = getChars(puzzle_word)
        amount = request.form['owmp_amount']
        amount = int(amount)

        showAnswers = False
        try:
            checkbox = request.form['owmp_checkbox']
            if checkbox == 'checked':
                showAnswers = True
        except:
            pass

        #print(f"logical chars: {type(langid.classify(puzzle_word))}")

        allPuzzles = []
        if langid.classify(puzzle_word)[0] == 'te':
            makePowerPoint('te', logicalWord, showAnswers, amount, type, size)
        else:
            makePowerPoint('en', logicalWord, showAnswers, amount, type, size)

        return send_file('C:/Users/bv2737dg/Documents/School/2022/499 Capstone (Wed)/Rebus/rebus_python/Rebus.pptx')
        #return render_template('oneWordManyPPT.html')
    else:
        return render_template('oneWordManyPPT.html')

@app.route('/manyWordsOnePuzzle', methods=['POST', 'GET'])
def manyWordsOnePuzzle():
    # return render_template('manyWordsOnePuzzle.html')
    if request.method == 'POST':
        puzzle_words = request.form['puzzle_words']
        wordList = puzzle_words.split() # turn the string into a list of words, splitting at the space

        allLogicalWords = [] # get a list of logical character lists
        for word in wordList:
            allLogicalWords.append(getChars(word))

        # loop and check every word for language
        # get a list of wordlists
        allPuzzles = []
        for i in range(len(wordList)):
            if langid.classify(wordList[i])[0] == 'te':
                allPuzzles.append(getManyLists(allLogicalWords[i], 'te', 1)[0])
            else:
                allPuzzles.append(getManyLists(allLogicalWords[i], 'en', 1)[0])
        return render_template('manyWordsOnePuzzle.html', load=True, puzzle_words=wordList, all_puzzles=allPuzzles)
    else:
        return render_template('manyWordsOnePuzzle.html', load=False)


@app.route('/manyWordsOnePuzzlePPT', methods=['POST', 'GET'])
def manyWordsOnePuzzlePPT():
    if request.method == 'POST':
        size = float(request.form['size_value'])
        type = request.form['image_size']
        puzzle_words = request.form['puzzle_words']
        wordList = puzzle_words.split()  # turn the string into a list of words, splitting at the space

        showAnswers = False
        try:
            checkbox = request.form['mwop_checkbox']
            if checkbox == 'checked':
                showAnswers = True
        except:
            pass

        allLogicalWords = []  # get a list of logical character lists
        for word in wordList:
            allLogicalWords.append(getChars(word))

        pr1 = Presentation()

        # no answers
        puzzles = []
        for i in range(len(wordList)):
            lw = allLogicalWords[i]
            n = 12
            chunks = [lw[i:i + n] for i in range(0, len(lw), n)]
            for chunk in chunks:
                if langid.classify(wordList[i])[0] == 'te':
                    puzzles.append(makeSlide(pr1, i + 1, 'te', chunk, False, size, type))
                else:
                    puzzles.append(makeSlide(pr1, i + 1, 'en', chunk, False, size, type))
        #show the answers this round
        if showAnswers:
            makeAnswerSlides(puzzles, pr1)

        pr1.save('Rebus.pptx')

        return send_file('C:/Users/bv2737dg/Documents/School/2022/499 Capstone (Wed)/Rebus/rebus_python/Rebus.pptx')
    else:
        return render_template('manyWordsOnePuzzlePPT.html', load=False)

def many_from_list(many_word_list):
    allPuzzles = []
    for word1 in many_word_list:
        wordStrings =''
        wordStrings += f'{"".join(word1)}: '
        #check_word1 = list(word1)
        tempList = many_word_list.copy()
        tempList.remove(word1)
        for letter in word1:
            string = ''.join(str(item) for item in tempList)
            if letter not in string:
                wordStrings += '??(not enough words to generate)'
                break
            for word2 in tempList:
                if letter in word2:
                    wordStrings += f'{word2.index(letter)+1}/{len(word2)}({"".join(word2)})  '
                    #print(f'{word2.index(letter)+1}/{len(word2)}({word2})  ', end='')
                    tempList.pop(tempList.index(word2))
                    break
        allPuzzles.append(wordStrings)
    return allPuzzles

@app.route('/manyFromList', methods=['POST', 'GET'])
def manyFromList():
    # return render_template('manyWordsOnePuzzle.html')
    if request.method == 'POST':
        puzzle_words = request.form['puzzle_words']
        wordList = puzzle_words.split() # turn the string into a list of words, splitting at the space

        allLogicalWords = [] # get a list of logical character lists
        for word in wordList:
            allLogicalWords.append(getChars(word))

        allPuzzles = many_from_list(allLogicalWords)

        return render_template('manyFromList.html', load=True, puzzle_words=puzzle_words, all_puzzles=allPuzzles)
    else:
        return render_template('manyFromList.html', load=False)

def makeManyFromListSlides(pr1, listOfWords):
    n = 8
    chunks = [listOfWords[i:i + n] for i in range(0, len(listOfWords), n)]
    for chunk in chunks:
        Layout = pr1.slide_layouts[6]
        slide = pr1.slides.add_slide(Layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
        textframe = textbox.text_frame
        for item in chunk:
            para = textframe.add_paragraph()
            para.text = str(item)
        textframe.fit_text()




@app.route('/manyFromListPPT', methods=['POST', 'GET'])
def manyFromListPPT():
    if request.method == 'POST':
        puzzle_words = request.form['puzzle_words']
        wordList = puzzle_words.split()  # turn the string into a list of words, splitting at the space

        allLogicalWords = []  # get a list of logical character lists
        for word in wordList:
            allLogicalWords.append(getChars(word))

        allPuzzles = many_from_list(allLogicalWords)

        pr2 = Presentation()
        makeManyFromListSlides(pr2, allPuzzles)
        pr2.save('manyFromList.pptx')

        return send_file('C:/Users/bv2737dg/Documents/School/2022/499 Capstone (Wed)/Rebus/rebus_python/manyFromList.pptx')
    else:
        return render_template('manyFromListPPT.html', load=False)

def one_from_given_list(puzzle_word, many_word_list):
    allPuzzles = []
    wordStrings =''
    wordStrings += f'{"".join(puzzle_word)}: '
    #check_word1 = list(word1)
    tempList = many_word_list.copy()
    for letter in puzzle_word:
        string = ''.join(str(item) for item in tempList)
        if letter not in string:
            wordStrings += '??(not enough words to generate)'
            break
        for word2 in tempList:
            if letter in word2:
                wordStrings += f'{word2.index(letter)+1}/{len(word2)}({"".join(word2)})  '
                #print(f'{word2.index(letter)+1}/{len(word2)}({word2})  ', end='')
                tempList.pop(tempList.index(word2))
                break
    allPuzzles.append(wordStrings)
    return allPuzzles

@app.route('/OneFromGivenList', methods=['POST', 'GET'])
def oneFromGivenList():
    # return render_template('manyWordsOnePuzzle.html')
    if request.method == 'POST':
        puzzle_word = request.form['puzzle_word']
        solution_words = request.form['solution_words']
        solution_words = solution_words.split(', ')

        allLogicalWords = [] # get a list of logical character lists
        for word in solution_words:
             allLogicalWords.append(getChars(word))

        logical_puzzle_word = getChars(puzzle_word)
        allPuzzles = one_from_given_list(logical_puzzle_word, allLogicalWords)
        # print(allPuzzles)
        return render_template('OneFromGivenList.html', load=True, puzzle_word=puzzle_word, all_puzzles=allPuzzles)
    else:
        return render_template('OneFromGivenList.html', load=False)


@app.route('/OneFromGivenListPPT', methods=['POST', 'GET'])
def oneFromGivenListPPT():
    # return render_template('manyWordsOnePuzzle.html')
    if request.method == 'POST':
        puzzle_word = request.form['puzzle_word']
        solution_words = request.form['solution_words']
        solution_words = solution_words.split(', ')

        allLogicalWords = [] # get a list of logical character lists
        for word in solution_words:
             allLogicalWords.append(getChars(word))

        allPuzzles = one_from_given_list(puzzle_word, allLogicalWords)

        pr2 = Presentation()
        makeManyFromListSlides(pr2, allPuzzles)
        pr2.save('OneFromGivenList.pptx')

        return send_file(
            'C:/Users/bv2737dg/Documents/School/2022/499 Capstone (Wed)/Rebus/rebus_python/OneFromGivenList.pptx')
    else:
        return render_template('OneFromGivenListPPT.html', load=False)

@app.route('/oneWithExclusion', methods=['POST', 'GET'])
def oneWithExclusion():
    if request.method == 'POST':
        puzzle_word = request.form['puzzle_word']
        exclusion_words = request.form['exclusion_words']
        if exclusion_words:
            exclusion_words = exclusion_words.split(', ')
        else:
            exclusion_words = []

        allLogicalSolutionWords = []  # get a list of logical character lists
        for word in exclusion_words:
            allLogicalSolutionWords.append(getChars(word))

        puzzle_word_logical = getChars(puzzle_word)
        # loop and check every word for language
        # get a list of wordlists
        allPuzzles = []
        if langid.classify(puzzle_word)[0] == 'te':
            allPuzzles.append(getManyLists(puzzle_word_logical, 'te', 1, exclusion_words)[0])
        else:
            allPuzzles.append(getManyLists(puzzle_word_logical, 'en', 1, exclusion_words)[0])

        return render_template('oneWithExclusion.html', load=True, puzzle_word=puzzle_word, all_puzzles=allPuzzles)
    else:
        return render_template('oneWithExclusion.html', load=False)


@app.route('/oneWithExclusionPPT', methods=['POST', 'GET'])
def oneWithExclusionPPT():
    # return render_template('manyWordsOnePuzzle.html')
    if request.method == 'POST':
        size = float(request.form['size_value'])
        type = request.form['image_size']
        showAnswers = False
        try:
            checkbox = request.form['owmp_checkbox']
            if checkbox == 'checked':
                showAnswers = True
        except:
            pass

        puzzle_word = request.form['puzzle_word']
        exclusion_words = request.form['exclusion_words']
        if exclusion_words:
            exclusion_words = exclusion_words.split(', ')
        else:
            exclusion_words = []

        allLogicalSolutionWords = []  # get a list of logical character lists
        for word in exclusion_words:
            allLogicalSolutionWords.append(getChars(word))

        pr1 = Presentation()
        logicalWord = getChars(puzzle_word)
        n = 12
        chunks = [logicalWord[i:i + n] for i in range(0, len(logicalWord), n)]

        puzzles = []
        for chunk in chunks:
            if langid.classify(puzzle_word)[0] == 'te':
               puzzles.append( makeSlide(pr1, 1, 'te', chunk, False, size, type, exclusion_words))
            else:
                puzzles.append(makeSlide(pr1, 1, 'en', chunk, False, size, type, exclusion_words))

        if showAnswers:
            makeAnswerSlides(puzzles, pr1)

        pr1.save('oneWithExclusion.pptx')

        return send_file(
            'C:/Users/bv2737dg/Documents/School/2022/499 Capstone (Wed)/Rebus/rebus_python/oneWithExclusion.pptx')
    else:
        return render_template('oneWithExclusionPPT.html', load=False)

if __name__ == "__main__":
    app.run(debug=True)

# def generatePPT():
#     with open('quotes_telugu.csv', encoding='utf-8') as file:
#         reader = csv.reader(file)
#         for row in reader:
#             line = row[1].replace("’", "'")
#             charList = getChars(line)
#




